"""
make_deck.py
Builds pitch_deck.pptx for the Artemis Analytics Quant Research Competition (Track 1).

A NON-TECHNICAL pitch deck telling an honest "null result" story:
  Act 1 - it looked like a winner (+63.6% / Sharpe 1.91)
  Act 2 - we tried to kill it (four-gate honesty scorecard -> 2/4 -> NOT DEPLOYABLE)

Run with the artemis conda env python:
    /Applications/anaconda3/envs/artemis/bin/python make_deck.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# --------------------------------------------------------------------------
# Paths
# --------------------------------------------------------------------------
HERE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(HERE, "figures")
OUT = os.path.join(HERE, "pitch_deck.pptx")

def fig(name):
    return os.path.join(FIG, name)

# Figure native aspect ratios (width / height)
ASPECT = {
    "fig1_equity_curve.png": 1.383,
    "fig2_quintiles.png":    1.659,
    "fig3_ic.png":           2.242,
    "fig4_oos.png":          2.032,
    "fig5_monthly.png":      2.242,
}

# --------------------------------------------------------------------------
# Palette  (ONE accent color: Artemis violet)
# --------------------------------------------------------------------------
ACCENT      = RGBColor(0x7C, 0x3A, 0xED)   # violet ~ rgb(124,58,237)
ACCENT_DK   = RGBColor(0x5B, 0x21, 0xB6)   # deeper violet for emphasis
INK         = RGBColor(0x1E, 0x1B, 0x2E)   # near-black dark text
SUBTLE      = RGBColor(0x55, 0x55, 0x66)   # muted gray text
BG          = RGBColor(0xFF, 0xFF, 0xFF)   # white slide background
BG_SOFT     = RGBColor(0xF6, 0xF4, 0xFC)   # very light violet wash for panels
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

PASS_FILL   = RGBColor(0xDD, 0xF2, 0xE3)   # soft green
PASS_TEXT   = RGBColor(0x14, 0x6B, 0x3A)
FAIL_FILL   = RGBColor(0xFB, 0xE0, 0xE0)   # soft red
FAIL_TEXT   = RGBColor(0x9B, 0x1C, 0x1C)
HDR_FILL    = ACCENT
HDR_TEXT    = WHITE
ROW_ALT     = RGBColor(0xF6, 0xF4, 0xFC)

FONT = "Calibri"

EMU_PER_IN = 914400
SW = Inches(13.333)
SH = Inches(7.5)

# --------------------------------------------------------------------------
# Presentation + blank layout
# --------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]   # truly blank


# --------------------------------------------------------------------------
# Helpers
# --------------------------------------------------------------------------
def add_slide():
    s = prs.slides.add_slide(BLANK)
    # paint white background
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    # send to back
    sp = bg._element
    sp.getparent().remove(sp)
    s.shapes._spTree.insert(2, sp)
    return s


def _set_font(run, size, bold=False, color=INK, name=FONT, italic=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = name
    run.font.color.rgb = color


def add_text(slide, left, top, width, height, lines,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    """lines: list of dicts -> {text,size,bold,color,space_after,italic,bullet,level,align}"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    for i, spec in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = spec.get("align", align)
        if spec.get("space_after") is not None:
            p.space_after = Pt(spec["space_after"])
        if spec.get("space_before") is not None:
            p.space_before = Pt(spec["space_before"])
        if spec.get("line_spacing") is not None:
            p.line_spacing = spec["line_spacing"]
        r = p.add_run()
        prefix = ""
        if spec.get("bullet"):
            prefix = "•   "
        r.text = prefix + spec["text"]
        _set_font(r, spec.get("size", 18), spec.get("bold", False),
                  spec.get("color", INK), spec.get("name", FONT),
                  spec.get("italic", False))
    return tb


def title_bar(slide, title, kicker=None):
    """Slim violet accent bar at top-left + title text + small kicker label."""
    # left accent block
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.0), Inches(0.0),
                                 Inches(0.22), Inches(1.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False

    top = Inches(0.42)
    if kicker:
        add_text(slide, Inches(0.55), Inches(0.30), Inches(11.5), Inches(0.35),
                 [{"text": kicker.upper(), "size": 13, "bold": True,
                   "color": ACCENT}])
        top = Inches(0.62)
    add_text(slide, Inches(0.55), top, Inches(12.2), Inches(1.0),
             [{"text": title, "size": 33, "bold": True, "color": INK}])
    # thin underline rule
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.58), Inches(1.52),
                                  Inches(2.1), Pt(3))
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()
    rule.shadow.inherit = False


def add_panel(slide, left, top, width, height, fill=BG_SOFT, line=None):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = fill
    if line is None:
        panel.line.fill.background()
    else:
        panel.line.color.rgb = line
        panel.line.width = Pt(1)
    panel.shadow.inherit = False
    # soften corner radius
    try:
        panel.adjustments[0] = 0.06
    except Exception:
        pass
    return panel


def add_picture_fit(slide, path, box_left, box_top, box_w, box_h,
                    align="center", valign="middle"):
    """Place picture scaled to fit a bounding box, preserving aspect ratio."""
    name = os.path.basename(path)
    ar = ASPECT.get(name)
    if ar is None:
        pic = slide.shapes.add_picture(path, box_left, box_top, width=box_w)
        return pic
    box_ar = box_w / box_h
    if ar >= box_ar:
        w = box_w
        h = int(box_w / ar)
    else:
        h = box_h
        w = int(box_h * ar)
    if align == "center":
        l = box_left + (box_w - w) // 2
    elif align == "left":
        l = box_left
    else:
        l = box_left + (box_w - w)
    if valign == "middle":
        t = box_top + (box_h - h) // 2
    elif valign == "top":
        t = box_top
    else:
        t = box_top + (box_h - h)
    pic = slide.shapes.add_picture(path, l, t, width=w, height=h)
    # subtle border line
    pic.line.color.rgb = RGBColor(0xDD, 0xD8, 0xEC)
    pic.line.width = Pt(0.75)
    return pic


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def stat_card(slide, left, top, w, h, big, label, big_color=ACCENT_DK,
              fill=BG_SOFT):
    add_panel(slide, left, top, w, h, fill=fill)
    add_text(slide, left, top + Inches(0.18), w, Inches(0.7),
             [{"text": big, "size": 30, "bold": True, "color": big_color}],
             align=PP_ALIGN.CENTER)
    add_text(slide, left, top + h - Inches(0.55), w, Inches(0.45),
             [{"text": label, "size": 13, "bold": True, "color": SUBTLE}],
             align=PP_ALIGN.CENTER)


# ==========================================================================
# SLIDE 1 — Title
# ==========================================================================
def slide_title():
    s = add_slide()
    # full-bleed violet band across the bottom third feel: use a left vertical
    # accent panel instead for a cleaner look
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), SH)
    band.fill.solid(); band.fill.fore_color.rgb = ACCENT
    band.line.fill.background(); band.shadow.inherit = False

    add_text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(0.5),
             [{"text": "ARTEMIS ANALYTICS QUANT RESEARCH COMPETITION  —  TRACK 1",
               "size": 15, "bold": True, "color": ACCENT}])

    add_text(s, Inches(0.85), Inches(1.95), Inches(11.8), Inches(2.0),
             [{"text": "Funding-Rate Mean Reversion:", "size": 46, "bold": True,
               "color": INK, "space_after": 2},
              {"text": "An Honest Null", "size": 46, "bold": True,
               "color": ACCENT_DK}])

    # thesis line
    add_panel(s, Inches(0.9), Inches(4.25), Inches(11.5), Inches(1.05),
              fill=BG_SOFT)
    add_text(s, Inches(1.2), Inches(4.42), Inches(10.9), Inches(0.8),
             [{"text": "A Sharpe-1.9 backtest that we deliberately tried to break "
                       "— and the honest reasons we would not trade it.",
               "size": 19, "bold": False, "italic": True, "color": INK,
               "line_spacing": 1.05}],
             anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.7),
             [{"text": "Crypto Factor Rebalancing  ·  Hyperliquid Perpetual Futures",
               "size": 15, "bold": True, "color": SUBTLE, "space_after": 2},
              {"text": "Presented by Daniel Connelly & Jack Brogan      ·      June 1, 2026",
               "size": 14, "color": SUBTLE}])

    set_notes(s, "Our entry is an honest null result. We found a funding-rate "
                 "mean-reversion strategy with a Sharpe of 1.9 in backtest, then "
                 "ran a battery of honesty tests designed to detect luck. It "
                 "failed two of four gates. We are presenting the diagnosis, not "
                 "a trade. The competition rewards the best thinking, not the "
                 "best P&L — that is exactly the spirit of this deck.")


# ==========================================================================
# SLIDE 2 — The idea in plain English
# ==========================================================================
def slide_idea():
    s = add_slide()
    title_bar(s, "The Idea, in Plain English", kicker="The intuition")

    add_text(s, Inches(0.6), Inches(1.85), Inches(7.0), Inches(0.9),
             [{"text": "Funding rate = a crowding thermometer.",
               "size": 24, "bold": True, "color": ACCENT_DK}])

    add_text(s, Inches(0.6), Inches(2.65), Inches(7.0), Inches(3.6),
             [{"text": "A perpetual future is a crypto bet with no expiry date. "
                       "To keep its price glued to the real (spot) price, traders "
                       "pay each other a small recurring fee called the funding "
                       "rate.", "size": 19, "color": INK, "space_after": 12,
               "line_spacing": 1.08},
              {"text": "When a trade gets crowded, that fee spikes:",
               "size": 19, "bold": True, "color": INK, "space_after": 8},
              {"text": "Everyone piling in LONG  ->  very high funding (longs pay).",
               "size": 18, "color": INK, "bullet": True, "space_after": 6,
               "line_spacing": 1.05},
              {"text": "Everyone piling in SHORT  ->  very negative funding "
                       "(shorts pay).", "size": 18, "color": INK, "bullet": True,
               "space_after": 12, "line_spacing": 1.05},
              {"text": "Our bet: the over-loved, over-leveraged side tends to "
                       "snap back. Crowded longs underperform; crowded shorts "
                       "outperform.", "size": 19, "bold": True, "color": ACCENT_DK,
               "line_spacing": 1.08}])

    # analogy panel on the right
    add_panel(s, Inches(8.1), Inches(1.95), Inches(4.7), Inches(4.6),
              fill=BG_SOFT)
    add_text(s, Inches(8.45), Inches(2.25), Inches(4.05), Inches(0.6),
             [{"text": "The everyday analogy", "size": 17, "bold": True,
               "color": ACCENT}])
    add_text(s, Inches(8.45), Inches(2.95), Inches(4.05), Inches(3.4),
             [{"text": "“Betting Against the Crowd.”",
               "size": 20, "bold": True, "italic": True, "color": INK,
               "space_after": 12, "line_spacing": 1.05},
              {"text": "When a stock — or a coin — becomes everyone's favorite, "
                       "it has often already had its run. The most loved end of "
                       "the market tends to earn less than its risk would "
                       "suggest.", "size": 17, "color": INK, "space_after": 12,
               "line_spacing": 1.12},
              {"text": "Funding rate just gives us a clean, real-money gauge of "
                       "exactly how crowded each coin is, every single day.",
               "size": 17, "color": INK, "line_spacing": 1.12}])

    set_notes(s, "Keep this plain. A perpetual future has no expiry, so an "
                 "exchange uses the funding rate — a real fee one side pays the "
                 "other — to keep it pinned to spot. Extreme funding means a "
                 "crowded, over-leveraged trade. Our hypothesis is mean "
                 "reversion: the crowded side snaps back. It's the same instinct "
                 "as betting against the most-hyped stock.")


# ==========================================================================
# SLIDE 3 — The trade
# ==========================================================================
def slide_trade():
    s = add_slide()
    title_bar(s, "The Trade", kicker="How it works")

    # four mechanic cards across the top
    cards = [
        ("Rank", "Every day, score the top-35 coins by how extreme their "
                 "funding is."),
        ("Buy the cheap side", "Go LONG the coins with the lowest "
                               "(most negative) funding — the crowded shorts."),
        ("Sell the hot side", "Go SHORT the coins with the highest funding — "
                              "the crowded longs."),
        ("Stay neutral", "Equal dollars long and short, so we bet on the gap, "
                         "not on the market direction."),
    ]
    cw = Inches(2.95)
    gap = Inches(0.16)
    left0 = Inches(0.6)
    top = Inches(2.0)
    ch = Inches(2.55)
    for i, (head, body) in enumerate(cards):
        l = Emu(int(left0) + i * (int(cw) + int(gap)))
        add_panel(s, l, top, cw, ch, fill=BG_SOFT)
        # number chip
        chip = s.shapes.add_shape(MSO_SHAPE.OVAL, l + Inches(0.25),
                                  top + Inches(0.25), Inches(0.55), Inches(0.55))
        chip.fill.solid(); chip.fill.fore_color.rgb = ACCENT
        chip.line.fill.background(); chip.shadow.inherit = False
        ctf = chip.text_frame; ctf.word_wrap = False
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run(); cr.text = str(i + 1)
        _set_font(cr, 20, True, WHITE)
        add_text(s, l + Inches(0.25), top + Inches(0.95), cw - Inches(0.5),
                 Inches(0.55),
                 [{"text": head, "size": 17, "bold": True, "color": ACCENT_DK}])
        add_text(s, l + Inches(0.25), top + Inches(1.45), cw - Inches(0.5),
                 Inches(1.0),
                 [{"text": body, "size": 14.5, "color": INK, "line_spacing": 1.08}])

    # definition strip
    add_panel(s, Inches(0.6), Inches(4.85), Inches(12.15), Inches(1.85),
              fill=RGBColor(0xF1, 0xEC, 0xFB))
    add_text(s, Inches(0.95), Inches(5.05), Inches(11.5), Inches(0.5),
             [{"text": "Two terms, defined once:", "size": 16, "bold": True,
               "color": ACCENT}])
    add_text(s, Inches(0.95), Inches(5.5), Inches(11.5), Inches(1.1),
             [{"text": "Market-neutral  —  equal money on both sides, so a broad "
                       "crypto crash or rally roughly cancels out. We only profit "
                       "from the spread between the two baskets.",
               "size": 16.5, "color": INK, "bullet": True, "space_after": 6,
               "line_spacing": 1.06},
              {"text": "Rebalanced daily  —  we refresh the baskets each day, "
                       "holding each position about two weeks (a 14-day horizon).",
               "size": 16.5, "color": INK, "bullet": True, "line_spacing": 1.06}])

    set_notes(s, "Mechanically simple. Each day we rank the top-35 coins by a "
                 "funding z-score, buy the lowest-funding names, short the "
                 "highest, equal-weight, dollar-neutral, rebalanced daily with "
                 "roughly a two-week holding horizon. Market-neutral means we are "
                 "betting on the spread, not on whether crypto goes up or down. "
                 "Note for later: the parameters were chosen by a sweep — we "
                 "mined them, which is the seed of the whole honesty story.")


# ==========================================================================
# SLIDE 4 — Act 1: it looked like a winner
# ==========================================================================
def slide_act1():
    s = add_slide()
    title_bar(s, "Act 1: It Looked Like a Winner", kicker="The backtest")

    # figure on the right
    add_picture_fit(s, fig("fig1_equity_curve.png"),
                    Inches(6.7), Inches(1.85), Inches(6.2), Inches(4.9))

    # stat cards on the left
    stat_card(s, Inches(0.6), Inches(2.0), Inches(2.85), Inches(1.5),
              "+63.6%", "Total return (~1 yr)")
    stat_card(s, Inches(3.6), Inches(2.0), Inches(2.85), Inches(1.5),
              "1.91", "Sharpe ratio")
    stat_card(s, Inches(0.6), Inches(3.65), Inches(2.85), Inches(1.5),
              "-19%", "Max drawdown", big_color=FAIL_TEXT)
    stat_card(s, Inches(3.6), Inches(3.65), Inches(2.85), Inches(1.5),
              "May '25–'26", "Backtest window")

    add_text(s, Inches(0.6), Inches(5.4), Inches(5.85), Inches(1.5),
             [{"text": "On paper, a strong result.", "size": 19, "bold": True,
               "color": ACCENT_DK, "space_after": 8},
              {"text": "A Sharpe near 2 with a manageable drawdown is the kind of "
                       "curve that gets a strategy funded. The honest question "
                       "is: is it real, or did we get lucky?",
               "size": 16, "color": INK, "line_spacing": 1.1}])

    set_notes(s, "This is the seductive part. Plus 63.6% over about a year, "
                 "Sharpe 1.91, worst drawdown only minus 19 percent. The shaded "
                 "region is out-of-sample and the dip is the January 2026 BTC "
                 "rally — hold that thought. A naive team ships this. We didn't, "
                 "because the parameters were mined and one year of crypto is "
                 "thin, regime-driven data.")


# ==========================================================================
# SLIDE 5 — Why we didn't trust it
# ==========================================================================
def slide_distrust():
    s = add_slide()
    title_bar(s, "Why We Didn't Trust It", kicker="Act 2 begins")

    add_text(s, Inches(0.6), Inches(1.9), Inches(12.0), Inches(0.8),
             [{"text": "A great-looking number is the beginning of the work, "
                       "not the end.", "size": 22, "bold": True, "color": INK}])

    reasons = [
        ("We mined the parameters",
         "The look-back windows and holding period were picked by searching "
         "over dozens of combinations. Search hard enough and something always "
         "looks great by chance."),
        ("Crypto history is short",
         "About one year of data. That isn't enough to tell skill from luck, "
         "especially for a daily-rebalanced strategy."),
        ("The market changes regime",
         "Crypto swings between calm, range-bound stretches and violent "
         "trends. A strategy can look brilliant in one regime and bleed in the "
         "next."),
    ]
    cw = Inches(3.93)
    gap = Inches(0.18)
    left0 = Inches(0.6)
    top = Inches(2.95)
    ch = Inches(3.05)
    for i, (head, body) in enumerate(reasons):
        l = Emu(int(left0) + i * (int(cw) + int(gap)))
        add_panel(s, l, top, cw, ch, fill=BG_SOFT)
        # accent top stripe
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, top, cw, Inches(0.12))
        stripe.fill.solid(); stripe.fill.fore_color.rgb = ACCENT
        stripe.line.fill.background(); stripe.shadow.inherit = False
        add_text(s, l + Inches(0.3), top + Inches(0.45), cw - Inches(0.6),
                 Inches(0.9),
                 [{"text": head, "size": 18, "bold": True, "color": ACCENT_DK,
                   "line_spacing": 1.0}])
        add_text(s, l + Inches(0.3), top + Inches(1.45), cw - Inches(0.6),
                 Inches(1.5),
                 [{"text": body, "size": 16, "color": INK, "line_spacing": 1.12}])

    add_text(s, Inches(0.6), Inches(6.25), Inches(12.0), Inches(0.7),
             [{"text": "So we stopped trying to confirm the strategy and started "
                       "trying to kill it.", "size": 18, "bold": True,
               "italic": True, "color": ACCENT_DK}])

    set_notes(s, "Three reasons for suspicion. One, we mined the parameters over "
                 "56 combinations — searching hard guarantees something looks good "
                 "by chance. Two, only about a year of data. Three, crypto is "
                 "regime-driven. The right scientific move is to invert the goal: "
                 "instead of confirming the strategy, actively try to falsify it.")


# ==========================================================================
# SLIDE 6 — The four-gate stress test (NATIVE TABLE)
# ==========================================================================
def slide_gates():
    s = add_slide()
    title_bar(s, "The Four-Gate Stress Test", kicker="Our honesty scorecard")

    add_text(s, Inches(0.6), Inches(1.7), Inches(12.0), Inches(0.5),
             [{"text": "Four pass/fail gates a strategy must clear before we would "
                       "ever deploy it.", "size": 17, "color": SUBTLE}])

    rows = 5
    cols = 3
    left = Inches(0.6)
    top = Inches(2.25)
    width = Inches(12.13)
    height = Inches(4.25)
    gtab = s.shapes.add_table(rows, cols, left, top, width, height).table

    # column widths
    gtab.columns[0].width = Inches(3.3)
    gtab.columns[1].width = Inches(6.43)
    gtab.columns[2].width = Inches(2.4)

    # turn off banding styling so our manual fills show cleanly
    tbl = gtab._tbl
    # remove the default table style for full manual control of fills
    tblPr = tbl.find(qn('a:tblPr'))
    if tblPr is not None:
        tblPr.set('firstRow', '0')
        tblPr.set('bandRow', '0')

    def cell_text(cell, text, size=15, bold=False, color=INK,
                  align=PP_ALIGN.LEFT):
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.15)
        cell.margin_right = Inches(0.12)
        cell.margin_top = Inches(0.06)
        cell.margin_bottom = Inches(0.06)
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        _set_font(r, size, bold, color)

    def fill_cell(cell, rgb):
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb

    # header
    headers = ["Gate", "The question it asks", "Result"]
    for c, h in enumerate(headers):
        cell = gtab.cell(0, c)
        fill_cell(cell, HDR_FILL)
        cell_text(cell, h, size=16, bold=True, color=HDR_TEXT,
                  align=(PP_ALIGN.CENTER if c == 2 else PP_ALIGN.LEFT))
    gtab.rows[0].height = Inches(0.6)

    # data rows: (gate, question, result, passed)
    data = [
        ("1. Statistical significance",
         "After correcting for how many combos we tried, is the edge still "
         "real? (0 of 56 survive; traded-spread p = 0.057)",
         "FAIL", False),
        ("2. Sign stability",
         "Does the edge point the same way on a sealed out-of-sample window? "
         "(+12.7% out-of-sample, ~76% of in-sample Sharpe)",
         "PASS", True),
        ("3. Overfitting penalty",
         "Discounting for how hard we searched, is it better than a coin "
         "flip? (Deflated Sharpe = 47%, bar is 95%)",
         "FAIL", False),
        ("4. Survives real costs",
         "Does it still work after realistic trading fees, even doubled? "
         "(gross 107.8% -> net 94.8% ann.; 82.7% at 2x)",
         "PASS", True),
    ]
    for i, (gate, q, res, passed) in enumerate(data, start=1):
        gtab.rows[i].height = Inches(0.91)
        c0 = gtab.cell(i, 0)
        fill_cell(c0, ROW_ALT if i % 2 == 0 else WHITE)
        cell_text(c0, gate, size=15.5, bold=True, color=INK)

        c1 = gtab.cell(i, 1)
        fill_cell(c1, ROW_ALT if i % 2 == 0 else WHITE)
        cell_text(c1, q, size=14, color=INK)

        c2 = gtab.cell(i, 2)
        fill_cell(c2, PASS_FILL if passed else FAIL_FILL)
        cell_text(c2, res, size=18, bold=True,
                  color=(PASS_TEXT if passed else FAIL_TEXT),
                  align=PP_ALIGN.CENTER)

    set_notes(s, "This scorecard is the heart of the entry. Four gates. Gate 1, "
                 "significance after multiple-testing correction: zero of 56 "
                 "combos survive, the traded spread is p = 0.057 — FAIL. Gate 2, "
                 "sign stability out-of-sample: +12.7%, about 76 percent of "
                 "in-sample Sharpe — PASS. Gate 3, the Deflated Sharpe penalizes "
                 "for search effort: 47 percent, a coin flip versus a 95 percent "
                 "bar — FAIL. Gate 4, costs: 107.8% gross to 94.8% net, +82.7% even at "
                 "double fees — PASS. Two of four.")


# ==========================================================================
# SLIDE 7 — Verdict
# ==========================================================================
def slide_verdict():
    s = add_slide()
    # full violet background for impact. NOTE: add_slide() already painted a white
    # background and sent it to z-index 2; this violet rect is left appended (in
    # front of that white rect) so it actually shows. The text boxes/pill added
    # below are appended afterwards and therefore sit in front of the violet.
    bgrect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bgrect.fill.solid(); bgrect.fill.fore_color.rgb = ACCENT_DK
    bgrect.line.fill.background(); bgrect.shadow.inherit = False

    add_text(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.6),
             [{"text": "THE VERDICT", "size": 18, "bold": True,
               "color": RGBColor(0xD6, 0xC6, 0xF7)}], align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.8), Inches(2.25), Inches(11.7), Inches(1.7),
             [{"text": "2 of 4 gates", "size": 64, "bold": True, "color": WHITE}],
             align=PP_ALIGN.CENTER)

    # NOT DEPLOYABLE pill
    pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(3.97), Inches(4.05), Inches(5.4),
                              Inches(1.05))
    pill.fill.solid(); pill.fill.fore_color.rgb = WHITE
    pill.line.fill.background(); pill.shadow.inherit = False
    try:
        pill.adjustments[0] = 0.5
    except Exception:
        pass
    ptf = pill.text_frame; ptf.word_wrap = False
    ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
    pp = ptf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    pr = pp.add_run(); pr.text = "NOT DEPLOYABLE"
    _set_font(pr, 30, True, FAIL_TEXT)

    add_text(s, Inches(1.8), Inches(5.6), Inches(9.7), Inches(1.2),
             [{"text": "It makes money and survives costs — but it cannot clear "
                       "an honest significance bar. That conclusion is the "
                       "deliverable.", "size": 19, "italic": True,
               "color": RGBColor(0xEA, 0xE2, 0xFB), "line_spacing": 1.15}],
             align=PP_ALIGN.CENTER)

    set_notes(s, "Two of four. The verdict is NOT DEPLOYABLE. We want to be very "
                 "clear: this is not us hiding a loser. It is us refusing to ship "
                 "a winner we cannot statistically defend. The discipline to say "
                 "no here is the point of the whole exercise.")


# ==========================================================================
# SLIDE 8 — The twist / key insight
# ==========================================================================
def slide_twist():
    s = add_slide()
    title_bar(s, "The Twist: How It Failed Matters", kicker="The key insight")

    add_text(s, Inches(0.6), Inches(1.8), Inches(6.0), Inches(0.6),
             [{"text": "It didn't fail by losing money.", "size": 22,
               "bold": True, "color": ACCENT_DK}])

    add_text(s, Inches(0.6), Inches(2.55), Inches(6.0), Inches(2.2),
             [{"text": "The strategy is big, survives realistic costs, and even "
                       "stayed positive on data it had never seen (+12.7% "
                       "out-of-sample).", "size": 17, "color": INK,
               "space_after": 10, "line_spacing": 1.12},
              {"text": "It died on exactly — and only — the two tests built to "
                       "detect luck:", "size": 17, "bold": True, "color": INK,
               "space_after": 6, "line_spacing": 1.1},
              {"text": "significance after a hard search, and the "
                       "search-penalized Deflated Sharpe.", "size": 17,
               "color": INK, "bullet": True, "line_spacing": 1.1}])

    # the fingerprint callout
    add_panel(s, Inches(0.6), Inches(5.05), Inches(5.95), Inches(1.75),
              fill=RGBColor(0xF1, 0xEC, 0xFB))
    add_text(s, Inches(0.9), Inches(5.25), Inches(5.4), Inches(1.4),
             [{"text": "A big, cost-robust, out-of-sample-positive number that "
                       "still can't clear an honest significance bar is the "
                       "textbook fingerprint of data-mining.",
               "size": 16.5, "bold": True, "italic": True, "color": ACCENT_DK,
               "line_spacing": 1.12}], anchor=MSO_ANCHOR.MIDDLE)

    # figure: OOS
    add_picture_fit(s, fig("fig4_oos.png"),
                    Inches(6.85), Inches(2.4), Inches(6.0), Inches(3.6))
    add_text(s, Inches(6.85), Inches(6.15), Inches(6.0), Inches(0.7),
             [{"text": "In-sample vs. out-of-sample: the edge carried forward "
                       "(+12.7%) — yet still fails the luck tests.",
               "size": 13.5, "italic": True, "color": SUBTLE,
               "line_spacing": 1.05}], align=PP_ALIGN.CENTER)

    set_notes(s, "Here is the insight the judges should remember. This strategy "
                 "fails in the most instructive way possible. It is not a "
                 "money-loser; it makes a lot, beats costs, and even carries into "
                 "out-of-sample. It dies specifically on the two tests for luck. "
                 "That precise signature — large, robust, OOS-positive, but "
                 "significance-failing — is the textbook fingerprint of "
                 "data-mining. Diagnosing that is our actual deliverable.")


# ==========================================================================
# SLIDE 9 — What we learned
# ==========================================================================
def slide_learned():
    s = add_slide()
    title_bar(s, "What We Learned", kicker="The lessons")

    # left: regime fragility w/ monthly fig
    add_text(s, Inches(0.6), Inches(1.8), Inches(6.0), Inches(0.55),
             [{"text": "1.  Regime is everything", "size": 19, "bold": True,
               "color": ACCENT_DK}])
    add_picture_fit(s, fig("fig5_monthly.png"),
                    Inches(0.6), Inches(2.35), Inches(6.05), Inches(2.6),
                    valign="top")
    add_text(s, Inches(0.6), Inches(5.1), Inches(6.05), Inches(1.8),
             [{"text": "The factor works when the market is range-bound and bled "
                       "during the January-2026 BTC rally — when crowded longs "
                       "kept being right, mean reversion is the wrong bet.",
               "size": 15.5, "color": INK, "line_spacing": 1.12}])

    # right: non-monotone tail w/ quintiles fig
    add_text(s, Inches(6.95), Inches(1.8), Inches(6.0), Inches(0.55),
             [{"text": "2.  It's a tail effect, not a clean factor",
               "size": 19, "bold": True, "color": ACCENT_DK}])
    add_picture_fit(s, fig("fig2_quintiles.png"),
                    Inches(6.95), Inches(2.35), Inches(5.9), Inches(2.6),
                    valign="top")
    add_text(s, Inches(6.95), Inches(5.1), Inches(5.9), Inches(1.8),
             [{"text": "Returns aren't a clean line across funding levels — the "
                       "relationship is non-monotone (the second-lowest-funding "
                       "bucket is anomalously weak). The edge lives only in the "
                       "extreme tails, fragile and easy to mis-trade.", "size": 15.5, "color": INK,
               "line_spacing": 1.12}])

    set_notes(s, "Three lessons. One, regime is everything — the factor works "
                 "range-bound and bled during the Jan-2026 BTC rally, because "
                 "when crowded longs keep being right, betting on reversion "
                 "loses. Two, the signal is a non-monotone tail effect, not a "
                 "clean linear factor — Q2 is anomalously weak (nearly as low "
                 "as Q5), so the edge lives only in the extreme "
                 "tails. Three, and most important, the rigor changed the "
                 "conclusion: a naive process would have shipped a Sharpe-1.9 "
                 "winner.")


# ==========================================================================
# SLIDE 10 — Why this is the right answer for Artemis + next steps
# ==========================================================================
def slide_why_right():
    s = add_slide()
    title_bar(s, "Why This Is the Right Answer", kicker="For this competition")

    # quote panel
    add_panel(s, Inches(0.6), Inches(1.85), Inches(12.15), Inches(1.6),
              fill=BG_SOFT)
    add_text(s, Inches(1.0), Inches(2.05), Inches(11.4), Inches(1.25),
             [{"text": "“We are not looking for the most profitable strategy "
                       "— we are looking for the best thinking. A negative-Sharpe "
                       "strategy that is correctly understood and honestly "
                       "evaluated will score better than impressive returns with "
                       "no critical analysis.”",
               "size": 17, "italic": True, "bold": True, "color": ACCENT_DK,
               "line_spacing": 1.12}], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.0), Inches(3.45), Inches(11.4), Inches(0.4),
             [{"text": "— Artemis competition brief", "size": 13, "bold": True,
               "color": SUBTLE}])

    # left: how we deliver on it
    add_text(s, Inches(0.6), Inches(4.1), Inches(6.0), Inches(0.5),
             [{"text": "We delivered the thinking", "size": 18, "bold": True,
               "color": INK}])
    add_text(s, Inches(0.6), Inches(4.65), Inches(6.0), Inches(2.4),
             [{"text": "A clear, plain-English hypothesis with real economic "
                       "logic.", "size": 15.5, "color": INK, "bullet": True,
               "space_after": 7, "line_spacing": 1.08},
              {"text": "A pre-committed pass/fail framework, not after-the-fact "
                       "rationalizing.", "size": 15.5, "color": INK,
               "bullet": True, "space_after": 7, "line_spacing": 1.08},
              {"text": "A correct diagnosis: an attractive backtest that is most "
                       "likely data-mining.", "size": 15.5, "color": INK,
               "bullet": True, "space_after": 7, "line_spacing": 1.08},
              {"text": "Named limitations and a concrete path forward.",
               "size": 15.5, "color": INK, "bullet": True, "line_spacing": 1.08}])

    # right: what would make it deployable
    add_panel(s, Inches(6.85), Inches(4.1), Inches(5.9), Inches(2.95),
              fill=RGBColor(0xF1, 0xEC, 0xFB))
    add_text(s, Inches(7.15), Inches(4.3), Inches(5.3), Inches(0.5),
             [{"text": "What would make it deployable", "size": 17, "bold": True,
               "color": ACCENT}])
    add_text(s, Inches(7.15), Inches(4.85), Inches(5.3), Inches(2.1),
             [{"text": "Add a regime filter — only trade when BTC is range-bound.",
               "size": 15, "color": INK, "bullet": True, "space_after": 6,
               "line_spacing": 1.06},
              {"text": "Get a longer, survivorship-corrected history.",
               "size": 15, "color": INK, "bullet": True, "space_after": 6,
               "line_spacing": 1.06},
              {"text": "Pair it with an offsetting trend / momentum factor.",
               "size": 15, "color": INK, "bullet": True, "space_after": 6,
               "line_spacing": 1.06},
              {"text": "Trade only the extreme tails, where the real edge lives.",
               "size": 15, "color": INK, "bullet": True, "line_spacing": 1.06}])

    set_notes(s, "Tie it back to the brief: best thinking over best backtest. We "
                 "delivered a clear hypothesis, a pre-committed pass/fail "
                 "framework, the correct diagnosis of data-mining, and named "
                 "limitations with a path forward. To make it real we would add a "
                 "regime filter, get longer survivorship-corrected data, pair it "
                 "with a trend factor, and trade only the tails.")


# ==========================================================================
# SLIDE 11 — Appendix: methodology + key numbers + reproducibility
# ==========================================================================
def slide_appendix():
    s = add_slide()
    title_bar(s, "Appendix: Methodology & Numbers", kicker="Reference")

    # left: methodology one-pager
    add_text(s, Inches(0.6), Inches(1.8), Inches(6.0), Inches(0.5),
             [{"text": "Methodology, one page", "size": 18, "bold": True,
               "color": ACCENT_DK}])
    add_text(s, Inches(0.6), Inches(2.35), Inches(6.1), Inches(4.6),
             [{"text": "Universe: top-35 Hyperliquid perpetuals by volume.",
               "size": 14.5, "color": INK, "bullet": True, "space_after": 6,
               "line_spacing": 1.05},
              {"text": "Signal: funding z = (21-day mean − 90-day mean) / "
                       "90-day std.", "size": 14.5, "color": INK, "bullet": True,
               "space_after": 6, "line_spacing": 1.05},
              {"text": "Trade: long lowest-z, short highest-z; equal-weight, "
                       "dollar-neutral.", "size": 14.5, "color": INK,
               "bullet": True, "space_after": 6, "line_spacing": 1.05},
              {"text": "Horizon: 14-day holding; rebalanced daily.",
               "size": 14.5, "color": INK, "bullet": True, "space_after": 6,
               "line_spacing": 1.05},
              {"text": "Params chosen by a multi-parameter sweep (hence the "
                       "overfitting tests).", "size": 14.5, "color": INK,
               "bullet": True, "space_after": 10, "line_spacing": 1.05},
              {"text": "Known limitations:", "size": 14.5, "bold": True,
               "color": INK, "space_after": 4},
              {"text": "survivorship bias (delisted coins silently dropped); "
                       "point-in-time universe bias (7 coins listed mid-study); "
                       "~1-year history; partially-burned out-of-sample.",
               "size": 13.5, "color": SUBTLE, "line_spacing": 1.1}])

    # right top: key numbers table (native)
    rows, cols = 8, 2
    left = Inches(7.0); top = Inches(2.35)
    width = Inches(5.75); height = Inches(3.55)
    kt = s.shapes.add_table(rows, cols, left, top, width, height).table
    kt.columns[0].width = Inches(3.75)
    kt.columns[1].width = Inches(2.0)
    tblPr = kt._tbl.find(qn('a:tblPr'))
    if tblPr is not None:
        tblPr.set('firstRow', '0'); tblPr.set('bandRow', '0')

    def kcell(cell, text, size=13.5, bold=False, color=INK,
              align=PP_ALIGN.LEFT, fillrgb=None):
        if fillrgb is not None:
            cell.fill.solid(); cell.fill.fore_color.rgb = fillrgb
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.1)
        cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
        p = cell.text_frame.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text
        _set_font(r, size, bold, color)

    kdata = [
        ("Metric", "Value", True),
        ("Total return (backtest)", "+63.6%", False),
        ("Sharpe ratio", "1.91", False),
        ("Max drawdown", "-19%", False),
        ("Out-of-sample return", "+12.7%", False),
        ("Bootstrap p-value (spread)", "0.057", False),
        ("Combos surviving correction", "0 of 56", False),
        ("Deflated Sharpe", "47%", False),
    ]
    kt.rows[0].height = Inches(0.42)
    for i, (k, v, hdr) in enumerate(kdata):
        if not hdr:
            kt.rows[i].height = Inches(0.44)
        if hdr:
            kcell(kt.cell(i, 0), k, size=14, bold=True, color=HDR_TEXT,
                  fillrgb=HDR_FILL)
            kcell(kt.cell(i, 1), v, size=14, bold=True, color=HDR_TEXT,
                  align=PP_ALIGN.CENTER, fillrgb=HDR_FILL)
        else:
            alt = ROW_ALT if i % 2 == 0 else WHITE
            kcell(kt.cell(i, 0), k, fillrgb=alt)
            kcell(kt.cell(i, 1), v, bold=True, color=ACCENT_DK,
                  align=PP_ALIGN.CENTER, fillrgb=alt)

    # right bottom: net-of-cost + reproducibility note
    add_panel(s, Inches(7.0), Inches(6.05), Inches(5.75), Inches(1.0),
              fill=RGBColor(0xF1, 0xEC, 0xFB))
    add_text(s, Inches(7.25), Inches(6.15), Inches(5.3), Inches(0.85),
             [{"text": "Costs ~6.5%/yr: 107.8% gross  ->  +94.8% net ann. "
                       "(+82.7% at 2x).", "size": 13, "bold": True,
               "color": ACCENT_DK, "space_after": 3, "line_spacing": 1.05},
              {"text": "Reproducible: code + notebooks on GitHub; run end-to-end "
                       "to regenerate every figure and number.",
               "size": 12, "italic": True, "color": SUBTLE, "line_spacing": 1.05}])

    set_notes(s, "Reference slide for Q&A. Methodology recap on the left, the "
                 "headline numbers in the table, and the cost / reproducibility "
                 "note at the bottom. Everything here is reproducible from the "
                 "GitHub repo and notebooks. Be ready to talk through the named "
                 "limitations: survivorship bias, point-in-time universe bias, "
                 "short history, and a partially-burned out-of-sample window.")


# --------------------------------------------------------------------------
# Build
# --------------------------------------------------------------------------
def build():
    slide_title()
    slide_idea()
    slide_trade()
    slide_act1()
    slide_distrust()
    slide_gates()
    slide_verdict()
    slide_twist()
    slide_learned()
    slide_why_right()
    slide_appendix()
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build()
    size = os.path.getsize(path)
    print("Saved:", path)
    print("Slides:", len(prs.slides._sldIdLst))
    print("Size (bytes):", size)
